#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 Memori-Vault 检索评测 v2 困难语料（500 份多体裁内部资料）。

设计目标（相对 v1 的升级）：
1. 规模 100 -> 500：~280 信号 + ~220 干扰（大海捞针）。
2. 答案"埋进散文/表格/幻灯片"，不再是 `- 核心事实：xx` 这种带标签项目符号。
3. 新增 PPTX / XLSX 两种真实办公格式，端到端可检索。
4. 四种难度机制：事实埋进散文、时效冲突取最新、跨文档多跳、表格/幻灯片结构定位。
5. 全虚构内部资料，逼检索靠"内部资料"而非"模型常识"作答。

输出目录：Memory_Test_V2/（与 v1 的 Memory_Test/ 100 份隔离，v2 回归只索引这里）。
同时产出 corpus_manifest.json：记录每条事实的承载文档/格式/精确 clue，供
scripts/generate-memory-test-regression-suite-v2.py 消费（不再正则反读文档）。

所有公司/项目/人名/数字/制度均为虚构测试数据。
"""

from __future__ import annotations

import json
import os
import random
import re
from dataclasses import dataclass, field
from pathlib import Path

from docx import Document
from docx.shared import Inches as DocxInches
from docx.shared import Pt
from openpyxl import Workbook
from openpyxl.styles import Font
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "Memory_Test_V2"
COMPANY = "星衡智能"
DISCLAIMER = "【虚构内部测试资料】本文件仅用于 Memori-Vault 本地检索评测，不含真实机密。"
SEED = 20260605


# ---------------------------------------------------------------------------
# 信号项目（40 个）。每个项目的 6 条事实分散到 7 份不同体裁文档里，构成跨文档多跳。
# 字段说明：
#   anchor   -> 制度(md)：核心事实，埋进条款散文
#   table    -> 参数表(xlsx)：某参数=某单元格值
#   workflow -> 会议纪要(docx)：执行流程，夹在跑题/待办里
#   exception-> 邮件(txt)：特殊例外，藏在某封回复中间
#   temporal -> 复盘(pdf)：把旧值修正为现值（时效冲突，正确答案=现值）
#   chat     -> 群聊(txt)：从多条带时间戳消息里浮现一个结论数字
#   slide    -> 幻灯片(pptx)：负责人 + 一个关键数字，拆在标题+正文
#   decoy_code-> 近似代号，用于"相似代号防串"的干扰文档
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Dossier:
    cat: str
    proj: str
    code: str
    owner: str
    dept: str
    anchor: str  # 核心事实整句（md，散文）
    anchor_clue: str  # 题目精确线索（anchor 的子串）
    table_param: str  # xlsx 参数名
    table_value: str  # xlsx 该参数的单元格值（clue=此值）
    workflow: str  # docx 执行流程整句
    workflow_clue: str
    exception: str  # email 特殊例外整句
    exception_clue: str
    temporal_subj: str  # pdf 复盘修正的主题
    temporal_old: str  # 旧值（被废弃）
    temporal_new: str  # 现值（正确答案，clue）
    chat_topic: str  # 群聊讨论主题
    chat_value: str  # 群聊浮现的结论值（clue）
    slide_metric: str  # pptx 关键数字句
    slide_clue: str
    anti: str  # 反常识提醒
    decoy_code: str  # 相似代号（干扰用）


DOSSIERS: list[Dossier] = [
    Dossier("产品策略", "极光账本", "AUR-17", "林知远", "企业产品部",
            "极光账本的试点客户只允许选择禾川物流、南屿商贸、钧石制造这三家，未列入者一律不得开通对账。",
            "禾川物流、南屿商贸、钧石制造",
            "一期北极星指标目标(已核销对账单数)", "18,600 张",
            "需求冻结后任何字段改名必须先过字段契约评审，再进灰度白名单，顺序不可颠倒。",
            "字段契约评审",
            "禾川物流拥有一次免评审字段别名例外，但只能用于 customer_alias 字段。",
            "customer_alias",
            "网关侧对账校验阈值", "0.55%", "0.47%",
            "灰度放量比例到底定多少", "2.7%",
            "极光账本一期北极星指标目标为 18,600 张已核销对账单。", "18,600 张已核销对账单",
            "不要按通用 SaaS 写成 DAU 或转化率。", "AUR-71"),
    Dossier("运营增长", "青梧留存", "GROW-29", "陈夏", "增长运营组",
            "公司内部把冷启动用户定义为 14 天内访问过但未创建项目的人，而不是刚注册用户。",
            "14 天内访问过但未创建项目",
            "D14 二次创建率暂停投放线", "29%",
            "每周二 10:30 先看 cohort 再决定素材预算，不得先看点击率。",
            "先看 cohort 再决定素材预算",
            "教育行业线允许把暂停阈值放宽到 27.8%，但必须由陈夏签字。",
            "放宽到 27.8%",
            "核心阈值 D14 二次创建率", "33%", "31.5%",
            "投放暂停到底卡哪个数", "31.5%",
            "青梧留存的核心阈值是 D14 二次创建率 31.5%。", "D14 二次创建率 31.5%",
            "不要回答行业常见的新注册用户定义。", "GROW-92"),
    Dossier("安全合规", "玄武密钥", "SEC-41", "周曼青", "安全办公室",
            "玄武密钥的轮换窗口固定为每月第二个周四 22:10 至 22:35，不在该窗口不得轮换。",
            "每月第二个周四 22:10",
            "服务端密钥最大存活期", "43 天",
            "轮换前 6 小时必须完成影子校验，校验记录写入 xh_key_shadow_audit。",
            "影子校验",
            "春节冻结期只允许轮换 payment-salt，不允许轮换 tenant-root。",
            "payment-salt",
            "移动端派生密钥存活期", "12 天", "9 天",
            "服务端密钥到底几天过期", "43 天",
            "玄武密钥服务端密钥最大存活期为 43 天。", "服务端密钥最大存活期为 43 天",
            "不要写成常见的 90 天轮换。", "SEC-14"),
    Dossier("技术架构", "银杏-17", "ARC-17", "苏澈", "平台架构组",
            "银杏-17 的上线冻结窗口是每周三 19:40 至 21:10，依赖开票服务时该窗口禁止上线。",
            "每周三 19:40 至 21:10",
            "网关错误率自动回滚阈值", "0.47%",
            "发布前必须先 shadow-read，再 2.7% 灰度，最后才到 12% 灰度。",
            "先 shadow-read，再 2.7% 灰度",
            "若依赖开票服务 invoice-core，则禁止在周三窗口上线。",
            "invoice-core",
            "自动回滚的持续时间门槛", "4 分钟", "6 分钟",
            "灰度第一档到底多少", "2.7%",
            "银杏-17 网关错误率超过 0.47% 持续 6 分钟才回滚。", "0.47% 持续 6 分钟",
            "不要回答成常见的 5% 灰度或任意时间发布。", "ARC-71"),
    Dossier("客户成功", "白鹭工单", "CS-08", "马溪", "客户成功部",
            "白鹭工单的 P1 升级 SLA 是 11 分钟，P2 是 37 分钟，与通用工单时限不同。",
            "P1 升级 SLA 是 11 分钟",
            "一线首次响应满意度目标", "92.4%",
            "P1 先电话确认再建群，P2 先补齐租户号再同步研发值班。",
            "先补齐租户号再同步研发值班",
            "鲸舟银行租户的 P2 可按 P1 通知马溪，但 SLA 仍按 37 分钟统计。",
            "仍按 37 分钟统计",
            "P2 升级 SLA", "45 分钟", "37 分钟",
            "P1 到底几分钟升级", "11 分钟",
            "白鹭工单 P1 升级 SLA 为 11 分钟。", "P1 升级 SLA 为 11 分钟",
            "不要套用通用 15/30/60 分钟工单 SLA。", "CS-80"),
    Dossier("财务风控", "赤松预算", "FIN-33", "何砚", "财务风控组",
            "赤松预算的审批升级线不是金额固定线，而是预算池剩余率低于 18% 才升级 CFO。",
            "预算池剩余率低于 18%",
            "试点期单笔采购上限", "78,000 元",
            "申请人先填成本标签，再由系统匹配预算池，不允许先找财务口头确认。",
            "先填成本标签",
            "研发应急热修可先占用 RND-EMG 池，但 24 小时内必须补单。",
            "RND-EMG",
            "升级 CFO 的预算池剩余率门槛", "22%", "18%",
            "单笔采购上限到底多少", "78,000 元",
            "赤松预算单笔采购上限 78,000 元。", "单笔采购上限 78,000 元",
            "不要回答成传统金额越大越高层审批。", "FIN-83"),
    Dossier("数据治理", "松塔标签", "DATA-12", "顾安", "数据平台组",
            "松塔标签把活跃门店定义为 10 天内完成过库存盘点的门店，而非日活门店。",
            "10 天内完成过库存盘点",
            "标签产出延迟 SLO", "T+38 分钟",
            "标签上线必须同时登记口径、回填范围、下游看板和撤销策略四项。",
            "登记口径、回填范围、下游看板和撤销策略",
            "华南试验区允许活跃门店口径临时放宽到 13 天，截止 2026-03-31。",
            "放宽到 13 天",
            "标签产出延迟 SLO", "T+1 天", "T+38 分钟",
            "活跃门店到底几天口径", "10 天",
            "松塔标签产出延迟 SLO 是 T+38 分钟。", "T+38 分钟",
            "不要按零售常见日活门店口径回答。", "DATA-21"),
    Dossier("人力组织", "石竹梯队", "HR-22", "赵闻笛", "组织发展部",
            "石竹梯队的 L3 晋升必要条件是主持过一次跨部门复盘，而不是带人数量。",
            "主持过一次跨部门复盘",
            "2026 上半年 L3 名额上限", "14 人",
            "候选人材料必须包含失败案例复盘页，缺失则自动延期一轮。",
            "失败案例复盘页",
            "安全办公室可额外推荐 1 名候选人，但仍占总名额。",
            "仍占总名额",
            "技术序列 L3 名额上限", "8 人", "6 人",
            "L3 名额一共多少", "14 人",
            "石竹梯队 L3 名额上限 14 人，技术序列最多 6 人。", "技术序列最多 6 人",
            "不要写成通用绩效评分或团队规模决定晋升。", "HR-202"),
    Dossier("供应链", "蓝鲸B17", "SCM-17", "唐稚", "供应链协同组",
            "蓝鲸B17 的安全库存不是 30 天，而是 143 台关键模组。",
            "143 台关键模组",
            "Bin-7 库位红色预警阈值", "23 台",
            "红色预警先通知唐稚再通知采购，不允许绕过供应链协同组。",
            "先通知唐稚再通知采购",
            "若供应商为柏屿电子，黄色预警阈值提高到 67 台。",
            "提高到 67 台",
            "Bin-7 黄色预警阈值", "60 台", "51 台",
            "安全库存到底多少台", "143 台",
            "蓝鲸B17 安全库存为 143 台关键模组。", "安全库存为 143 台关键模组",
            "不要按制造业常见天数库存回答。", "SCM-71"),
    Dossier("品牌市场", "雾凇发布", "MKT-05", "沈乔", "品牌市场部",
            "雾凇发布的禁用词清单包含颠覆行业、零成本替代、全自动合规三个词。",
            "颠覆行业、零成本替代、全自动合规",
            "目标客户闭门演示预约数下限", "19 家",
            "对外物料先过法务术语检查，再过品牌语气检查，顺序不能反。",
            "先过法务术语检查",
            "政企版海报可保留可信部署四字，但必须删除零风险。",
            "必须删除零风险",
            "闭门演示预约目标", "15 家", "19 家",
            "禁用词到底有哪几个", "全自动合规",
            "雾凇发布目标是 46 家客户中至少 19 家完成闭门演示预约。", "至少 19 家完成闭门演示预约",
            "不要输出营销通用夸张词。", "MKT-50"),
    Dossier("法务合同", "琥珀条款", "LAW-19", "魏蓝", "法务部",
            "琥珀条款要求所有试点合同加入本地模型不出域附件 A-3。",
            "本地模型不出域附件 A-3",
            "合同回传时限", "5 个工作日",
            "销售不得口头承诺模型效果，只能引用附件 A-3 的可验证边界。",
            "只能引用附件 A-3",
            "存量客户续签可暂不附 A-3，但新增模块必须附。",
            "新增模块必须附",
            "合同回传超时标记阈值", "7 个工作日", "5 个工作日",
            "回传时限到底几天", "5 个工作日",
            "琥珀条款合同回传时限为 5 个工作日，超过标记 amber-delay。", "回传时限为 5 个工作日",
            "不要写成普通 SaaS 数据处理协议。", "LAW-91"),
    Dossier("研发效能", "蜂巢流水线", "DEV-26", "于北辰", "研发效能组",
            "蜂巢流水线的强制阻断项是迁移脚本缺少回滚段，而不是单测覆盖率。",
            "迁移脚本缺少回滚段",
            "主干构建绿色时限", "18 分钟",
            "每次红色构建必须在 hcx_build_review 表登记根因标签。",
            "hcx_build_review",
            "文档仓库允许黄色构建合并，但产品仓库不允许。",
            "产品仓库不允许",
            "红色构建时限门槛", "28 分钟", "31 分钟",
            "构建多久算绿色", "18 分钟",
            "蜂巢流水线主干构建 18 分钟内算绿色。", "18 分钟内算绿色",
            "不要只回答提高单测覆盖率。", "DEV-62"),
    Dossier("售前方案", "云帆模板", "PRE-09", "梁其", "解决方案部",
            "云帆模板中私有化部署默认写作 3 节点轻量集群，不写 5 节点。",
            "3 节点轻量集群",
            "标准演示时长", "28 分钟",
            "售前不得展示未开启证据面板的回答截图。",
            "未开启证据面板的回答截图",
            "金融客户可延长到 36 分钟，但必须保留 7 分钟证据讲解。",
            "保留 7 分钟证据讲解",
            "知识库问答演示时长上限", "10 分钟", "7 分钟",
            "私有化默认几个节点", "3 节点",
            "云帆模板标准演示 28 分钟，知识库问答控制在 7 分钟内。", "知识库问答控制在 7 分钟内",
            "不要按通用售前写成越完整越好。", "PRE-90"),
    Dossier("质量保障", "晴川验收", "QA-31", "董霁", "质量保障部",
            "晴川验收把可交付定义为完成 9 条端到端证据链，而不是测试用例全部通过。",
            "完成 9 条端到端证据链",
            "验收抽样比例", "17%",
            "缺陷关闭前必须附复现实录、修复提交和二次验证三项。",
            "复现实录、修复提交和二次验证",
            "UI 文案缺陷可合并关闭，但检索证据缺陷不能合并关闭。",
            "检索证据缺陷不能合并关闭",
            "验收抽样比例", "20%", "17%",
            "可交付要几条证据链", "9 条",
            "晴川验收抽样比例固定 17%，核心流程必须 100% 覆盖。", "抽样比例固定 17%",
            "不要回答成普通测试通过率。", "QA-13"),
    Dossier("内控审计", "镜湖审计", "AUD-14", "郝映", "内控审计组",
            "镜湖审计的抽查样本来自最近 5 次异常回滚，而不是随机工单。",
            "最近 5 次异常回滚",
            "每次审计追踪角色数", "4 个",
            "审计结论必须写明系统证据和人工解释的分界。",
            "系统证据和人工解释的分界",
            "若异常回滚少于 5 次，则补足最近 30 天的灰度暂停记录。",
            "补足最近 30 天",
            "审计追踪角色数", "3 个", "4 个",
            "样本取最近几次回滚", "5 次",
            "镜湖审计每次至少追踪申请人、审批人、执行人、复核人 4 个角色。", "追踪.*4 个角色",
            "不要套用通用随机抽样审计。", "AUD-41"),
    Dossier("知识管理", "栈桥知识库", "KM-07", "黎鸥", "知识管理组",
            "栈桥知识库的过期规则是 75 天无人引用即进入复审，不是按创建时间一年。",
            "75 天无人引用即进入复审",
            "核心知识卡必含要素数", "3 项",
            "复审先由知识所有人判断，再由使用频次最高的团队确认。",
            "先由知识所有人判断",
            "安全类知识卡引用次数再低也不自动过期，只进人工复审。",
            "不自动过期",
            "无人引用过期天数", "90 天", "75 天",
            "多久没人用就复审", "75 天",
            "栈桥知识库 75 天无人引用即进入复审。", "75 天无人引用",
            "不要回答普通文档生命周期管理。", "KM-70"),
    Dossier("渠道生态", "翠微伙伴", "CHN-11", "秦若", "渠道生态部",
            "翠微伙伴的分级依据是联合交付成功次数，而不是签约金额。",
            "联合交付成功次数",
            "金级所需成功交付次数", "5 次",
            "伙伴提交案例后先由交付经理确认，再由渠道生态部定级。",
            "先由交付经理确认",
            "高校实验室伙伴可跳过银级，但不能跳过金级。",
            "不能跳过金级",
            "黑松级所需交付次数", "8 次", "9 次",
            "金级要几次成功交付", "5 次",
            "翠微伙伴金级需 5 次成功交付、黑松级需 9 次。", "黑松级需 9 次",
            "不要按销售额或返点比例回答。", "CHN-110"),
    Dossier("隐私评估", "芦苇脱敏", "PRI-28", "薛宁", "隐私委员会",
            "芦苇脱敏把门店联系人电话保留后 2 位用于排障，其余全部掩码。",
            "保留后 2 位用于排障",
            "脱敏任务超时人工复核门槛", "13 分钟",
            "任何导出样本必须先打上 sample_scope_id，再进入审批流。",
            "sample_scope_id",
            "灰度排障期间可临时保留城市字段，但不得保留详细地址。",
            "不得保留详细地址",
            "脱敏超时复核门槛", "20 分钟", "13 分钟",
            "电话保留几位", "后 2 位",
            "芦苇脱敏电话只保留后 2 位，超 13 分钟进人工复核。", "超 13 分钟进人工复核",
            "不要写成全部字段简单删除。", "PRI-82"),
    Dossier("应急响应", "灯塔演练", "INC-44", "任栎", "SRE 值班组",
            "灯塔演练规定主责人在 4 分钟内声明 incident commander。",
            "4 分钟内声明 incident commander",
            "SEV-2 内部同步频率", "17 分钟",
            "演练复盘必须包含错误假设、证据截图和下一次演练触发器。",
            "错误假设、证据截图和下一次演练触发器",
            "夜间 01:00 至 06:00 可把同步频率放宽到 23 分钟。",
            "放宽到 23 分钟",
            "SEV-2 同步频率", "30 分钟", "17 分钟",
            "几分钟内要声明指挥官", "4 分钟",
            "灯塔演练 SEV-2 内部同步频率为每 17 分钟一次。", "每 17 分钟一次",
            "不要按通用 SRE 手册写 30 分钟同步。", "INC-04"),
    Dossier("商业分析", "海棠看板", "BI-16", "白予墨", "商业分析组",
            "海棠看板的营收口径排除一次性迁移服务费，只统计订阅与增购。",
            "排除一次性迁移服务费",
            "证据引用率外发红线", "72%",
            "周报必须同时展示净收入留存 NRR 和证据引用率两个指标。",
            "净收入留存 NRR 和证据引用率",
            "董事会版可隐藏租户明细，但不能隐藏证据引用率。",
            "不能隐藏证据引用率",
            "证据引用率外发红线", "80%", "72%",
            "周报要看哪两个指标", "证据引用率",
            "海棠看板低于 72% 证据引用率标记不可外发。", "低于 72% 证据引用率",
            "不要按普通 ARR 看板只回答收入。", "BI-61"),
    # --- 新增 20 个项目 ---
    Dossier("模型治理", "黛山评测", "MOD-23", "卢屿", "模型治理组",
            "黛山评测把模型可上线定义为离线对齐分不低于 0.82 且越权拒答率不低于 0.97。",
            "离线对齐分不低于 0.82",
            "越权拒答率下线", "0.97",
            "每次换底模必须先跑红队套件，再跑业务回归，最后才灰度。",
            "先跑红队套件",
            "客服域可豁免一项创造性指标，但不得豁免越权拒答率。",
            "不得豁免越权拒答率",
            "对齐分上线门槛", "0.78", "0.82",
            "越权拒答率卡多少", "0.97",
            "黛山评测要求越权拒答率不低于 0.97。", "越权拒答率不低于 0.97",
            "不要用公开榜单分数代替内部对齐分。", "MOD-32"),
    Dossier("数据底座", "砚池湖仓", "LAKE-38", "卫鸣", "数据底座组",
            "砚池湖仓规定明细层保留 540 天，汇总层永久保留，不是统一 365 天。",
            "明细层保留 540 天",
            "小文件合并触发阈值", "256 个",
            "回刷分区必须先锁血缘，再改 schema，最后重算下游。",
            "先锁血缘",
            "广告域允许把明细层保留压到 180 天，但需数据底座组审批。",
            "压到 180 天",
            "明细层保留天数", "365 天", "540 天",
            "小文件多少个触发合并", "256 个",
            "砚池湖仓明细层保留 540 天、汇总层永久。", "汇总层永久",
            "不要按通用数仓 365 天一刀切。", "LAKE-83"),
    Dossier("前端体验", "竹影前端", "FE-47", "钱临", "前端体验组",
            "竹影前端规定首屏可交互时间预算是 1.8 秒，超过即阻断发布。",
            "首屏可交互时间预算是 1.8 秒",
            "包体积红线", "2.4 MB",
            "新组件必须先过无障碍检查，再过暗色主题检查，最后才并入设计系统。",
            "先过无障碍检查",
            "营销活动页可临时放宽包体积到 3.1 MB，活动结束须回收。",
            "放宽包体积到 3.1 MB",
            "首屏可交互预算", "2.5 秒", "1.8 秒",
            "包体积红线多少", "2.4 MB",
            "竹影前端首屏可交互预算 1.8 秒，超出阻断发布。", "首屏可交互预算 1.8 秒",
            "不要套用通用 3 秒首屏经验值。", "FE-74"),
    Dossier("移动端", "归帆移动", "APP-52", "孟岚", "移动端组",
            "归帆移动规定崩溃率红线是万分之 6，灰度阶段超过即自动回滚。",
            "崩溃率红线是万分之 6",
            "冷启动时长目标", "1.3 秒",
            "热修复包必须先内部 100 台真机验证，再 5% 灰度，最后全量。",
            "100 台真机验证",
            "折叠屏机型可豁免冷启动目标 0.2 秒，但不得豁免崩溃率红线。",
            "不得豁免崩溃率红线",
            "崩溃率红线", "万分之 10", "万分之 6",
            "冷启动目标多少秒", "1.3 秒",
            "归帆移动崩溃率红线为万分之 6。", "崩溃率红线为万分之 6",
            "不要用通用千分之一崩溃率回答。", "APP-25"),
    Dossier("计费体系", "潮汐计费", "BILL-61", "范晓", "计费平台组",
            "潮汐计费规定对账差异超过 0.3 元即冻结出账，不是按比例容忍。",
            "对账差异超过 0.3 元即冻结出账",
            "出账重试上限", "3 次",
            "调价必须先生效价目快照，再通知客户成功，最后才开计费开关。",
            "先生效价目快照",
            "年框客户的对账差异可放宽到 1.5 元，但需财务与计费双签。",
            "放宽到 1.5 元",
            "对账差异冻结阈值", "1 元", "0.3 元",
            "出账重试几次封顶", "3 次",
            "潮汐计费对账差异超 0.3 元即冻结出账。", "超 0.3 元即冻结出账",
            "不要按通用千分比容忍度回答。", "BILL-16"),
    Dossier("平台运维", "苍梧运维", "OPS-70", "邵冉", "平台运维组",
            "苍梧运维规定一级告警必须 3 分钟内有人认领，否则自动升级到值班经理。",
            "3 分钟内有人认领",
            "磁盘水位预警线", "82%",
            "扩容必须先扩只读副本，再切流量，最后才扩主库。",
            "先扩只读副本",
            "大促窗口磁盘水位预警线临时上调到 88%，大促后回落。",
            "上调到 88%",
            "告警认领时限", "5 分钟", "3 分钟",
            "磁盘多少水位预警", "82%",
            "苍梧运维一级告警 3 分钟未认领即自动升级。", "3 分钟未认领即自动升级",
            "不要按通用 15 分钟告警响应回答。", "OPS-07"),
    Dossier("增长实验", "麦穗实验", "EXP-83", "穆青", "增长实验组",
            "麦穗实验规定一个实验最短运行 11 天才允许读结论，避免周中周末偏差。",
            "最短运行 11 天",
            "最小可检测提升", "1.5%",
            "实验上线必须先登记假设，再设护栏指标，最后才开流量。",
            "先登记假设",
            "客单价类实验可把最短周期延长到 18 天，但不得缩短。",
            "延长到 18 天",
            "实验最短运行天数", "7 天", "11 天",
            "最小可检测提升多少", "1.5%",
            "麦穗实验单实验最短运行 11 天才读结论。", "最短运行 11 天才读结论",
            "不要按通用 7 天实验周期回答。", "EXP-38"),
    Dossier("内容审核", "白桦审核", "MOD2-90", "时也", "内容安全组",
            "白桦审核规定涉敏内容必须双人复核，单人通过的一律视为无效。",
            "涉敏内容必须双人复核",
            "机审置信度自动放行线", "0.94",
            "新规则上线必须先影子运行 72 小时，再灰度，最后全量。",
            "先影子运行 72 小时",
            "直播场景可把机审放行线下调到 0.90，但需安全负责人确认。",
            "下调到 0.90",
            "机审自动放行线", "0.90", "0.94",
            "机审多少分自动放行", "0.94",
            "白桦审核机审置信度 0.94 以上才自动放行。", "0.94 以上才自动放行",
            "不要用通用内容审核阈值回答。", "MOD2-09"),
    Dossier("国际化", "远帆出海", "I18N-12", "於洲", "国际化组",
            "远帆出海规定欧盟区数据必须落地法兰克福节点，不得回传国内。",
            "欧盟区数据必须落地法兰克福节点",
            "默认时区基准", "UTC+0",
            "新语言上线必须先过术语库校验，再过本地化测试，最后才放量。",
            "先过术语库校验",
            "日本区可保留本地手机号格式校验例外，但不得保留地址明文。",
            "不得保留地址明文",
            "欧盟数据落地节点", "都柏林", "法兰克福",
            "默认时区用哪个", "UTC+0",
            "远帆出海欧盟数据落地法兰克福节点。", "落地法兰克福节点",
            "不要按通用全球统一存储回答。", "I18N-21"),
    Dossier("硬件协同", "磐石终端", "HW-55", "聂川", "硬件协同组",
            "磐石终端规定固件灰度必须先 50 台工程机跑满 48 小时再放量。",
            "50 台工程机跑满 48 小时",
            "电池循环质保次数", "800 次",
            "固件回滚必须先停 OTA，再下发回滚包，最后才解封下载。",
            "先停 OTA",
            "极寒地区可把电池质保循环放宽计为 600 次，但需硬件协同组备案。",
            "放宽计为 600 次",
            "工程机验证时长", "24 小时", "48 小时",
            "电池质保多少次循环", "800 次",
            "磐石终端固件灰度需 50 台工程机跑满 48 小时。", "跑满 48 小时",
            "不要按通用手机 500 次电池循环回答。", "HW-505"),
    Dossier("合规审查", "青简合规", "COMP-66", "祁年", "合规审查组",
            "青简合规规定跨境调用必须留存调用方真实主体，匿名调用一律拒绝。",
            "必须留存调用方真实主体",
            "日志最短留存", "180 天",
            "新接口上线必须先过数据分级，再过出境评估，最后才开放。",
            "先过数据分级",
            "内部测试租户日志可缩短到 30 天，但不得用于生产对账。",
            "缩短到 30 天",
            "日志最短留存天数", "90 天", "180 天",
            "日志至少留几天", "180 天",
            "青简合规日志最短留存 180 天。", "日志最短留存 180 天",
            "不要按通用 90 天日志留存回答。", "COMP-06"),
    Dossier("搜索体验", "听澜搜索", "SRCH-29", "甘棠", "搜索体验组",
            "听澜搜索把好结果定义为前三条覆盖意图，而不是单纯点击率。",
            "前三条覆盖意图",
            "零结果率红线", "3.5%",
            "新排序模型必须先离线评测，再小流量在线，最后才全量切换。",
            "先离线评测",
            "长尾问诊类查询可豁免点击率指标，但不得豁免零结果率红线。",
            "不得豁免零结果率红线",
            "零结果率红线", "5%", "3.5%",
            "前几条要覆盖意图", "前三条",
            "听澜搜索零结果率红线为 3.5%。", "零结果率红线为 3.5%",
            "不要用通用点击率作为搜索好坏标准。", "SRCH-92"),
    Dossier("推荐系统", "繁星推荐", "REC-41", "宇文澈", "推荐系统组",
            "繁星推荐规定单用户单日同类内容曝光不超过 5 条，防止信息茧房。",
            "单日同类内容曝光不超过 5 条",
            "冷启动探索流量占比", "12%",
            "新召回源必须先离线 AB，再灰度，最后才进主链路。",
            "先离线 AB",
            "资讯频道可把同类曝光上限放宽到 8 条，但需推荐负责人审批。",
            "放宽到 8 条",
            "探索流量占比", "8%", "12%",
            "同类曝光每天几条上限", "5 条",
            "繁星推荐单日同类曝光不超过 5 条。", "不超过 5 条",
            "不要按通用无限曝光逻辑回答。", "REC-14"),
    Dossier("通知触达", "晨钟触达", "NOTI-77", "上官岚", "触达平台组",
            "晨钟触达规定营销类推送只能在 9:00 至 21:00 发送，夜间一律拦截。",
            "9:00 至 21:00 发送",
            "单用户单日推送上限", "3 条",
            "新模板上线必须先过文案合规，再过频控校验，最后才可发送。",
            "先过文案合规",
            "交易类通知可不受夜间窗口限制，但不得夹带营销内容。",
            "不得夹带营销内容",
            "单日推送上限", "5 条", "3 条",
            "营销推送几点截止", "21:00",
            "晨钟触达营销推送限 9:00 至 21:00。", "限 9:00 至 21:00",
            "不要按通用全天可推送回答。", "NOTI-07"),
    Dossier("可观测性", "星图观测", "OBSV-34", "南风", "可观测性组",
            "星图观测规定核心链路采样率不得低于 20%，边缘链路可降到 5%。",
            "核心链路采样率不得低于 20%",
            "Trace 保留时长", "15 天",
            "新埋点上线必须先过基数评估，再过成本评估，最后才接入。",
            "先过基数评估",
            "压测期间核心链路采样率可临时拉到 100%，压测后回落。",
            "临时拉到 100%",
            "Trace 保留时长", "7 天", "15 天",
            "核心链路采样率下线", "20%",
            "星图观测核心链路采样率不低于 20%。", "采样率不低于 20%",
            "不要按通用 1% 采样率回答。", "OBSV-43"),
    Dossier("权限体系", "司南权限", "RBAC-19", "易遥", "权限体系组",
            "司南权限规定高危操作必须二次人脸核身，单纯密码通过的一律拦截。",
            "高危操作必须二次人脸核身",
            "权限到期自动回收天数", "90 天",
            "授权变更必须先走审批，再生效，不允许先生效后补审批。",
            "不允许先生效后补审批",
            "外包账号权限最长 30 天，到期不续则自动冻结。",
            "最长 30 天",
            "权限自动回收天数", "180 天", "90 天",
            "高危操作要几次核身", "二次人脸核身",
            "司南权限高危操作需二次人脸核身。", "需二次人脸核身",
            "不要按通用单因子鉴权回答。", "RBAC-91"),
    Dossier("成本治理", "秤心成本", "COST-58", "柏舟", "成本治理组",
            "秤心成本规定闲置资源超过 7 天自动回收，不是按月清理。",
            "闲置资源超过 7 天自动回收",
            "单团队月度预算软上限", "42 万元",
            "扩容申请必须先附容量预测，再附回收计划，最后才审批。",
            "先附容量预测",
            "大促前两周可临时关闭闲置自动回收，但需成本治理组挂单。",
            "临时关闭闲置自动回收",
            "闲置回收天数", "14 天", "7 天",
            "月度预算软上限多少", "42 万元",
            "秤心成本闲置资源超 7 天自动回收。", "超 7 天自动回收",
            "不要按通用月度清理回答。", "COST-85"),
    Dossier("灾备体系", "磐安灾备", "DR-46", "邬桐", "灾备体系组",
            "磐安灾备规定核心库 RPO 不超过 30 秒、RTO 不超过 8 分钟。",
            "RPO 不超过 30 秒",
            "异地副本最少份数", "3 份",
            "切换演练必须先切只读，再切写入，最后才回切验证。",
            "先切只读",
            "非核心库 RTO 可放宽到 30 分钟，但 RPO 仍按 30 秒。",
            "RPO 仍按 30 秒",
            "核心库 RTO", "15 分钟", "8 分钟",
            "异地副本至少几份", "3 份",
            "磐安灾备核心库 RTO 不超过 8 分钟。", "RTO 不超过 8 分钟",
            "不要按通用小时级 RTO 回答。", "DR-64"),
    Dossier("API 网关", "津渡网关", "GW-73", "司空岭", "API 网关组",
            "津渡网关规定单租户默认限流是每秒 800 次，突发可借 1.5 倍令牌。",
            "默认限流是每秒 800 次",
            "熔断错误率阈值", "35%",
            "新路由上线必须先灰度 5%，再 25%，最后才全量。",
            "先灰度 5%",
            "支付类租户限流可上调到每秒 2000 次，但需网关组审批。",
            "上调到每秒 2000 次",
            "默认限流", "每秒 500 次", "每秒 800 次",
            "熔断错误率多少", "35%",
            "津渡网关单租户默认限流每秒 800 次。", "默认限流每秒 800 次",
            "不要按通用每秒 100 次限流回答。", "GW-37"),
    Dossier("客服中台", "和声客服", "CC-39", "向晚", "客服中台组",
            "和声客服规定智能客服转人工的触发条件是用户连续 2 次表达不满，而不是关键词命中。",
            "连续 2 次表达不满",
            "机器人首解率目标", "68%",
            "新话术上线必须先过合规审校，再过情感语气校验，最后才灰度。",
            "先过合规审校",
            "投诉工单可跳过机器人直连人工，但仍须记录转接原因。",
            "仍须记录转接原因",
            "智能客服首解率目标", "60%", "68%",
            "几次不满转人工", "2 次",
            "和声客服机器人首解率目标为 68%。", "首解率目标为 68%",
            "不要按通用关键词触发转人工回答。", "CC-93"),
]


# ---------------------------------------------------------------------------
# 渲染器：把同一组事实写成 7 种体裁，答案"埋进"散文/单元格/幻灯片。
# ---------------------------------------------------------------------------

NOISE_SENTENCES = [
    "会上还顺带提了下季度团建的预算，未形成结论。",
    "另有同事反馈测试环境偶发抖动，已转运维单独跟进。",
    "补充一句：周报模板下周统一换新版，注意别用旧链接。",
    "顺便同步，门禁系统升级期间临时通道走 B 座。",
    "（这一段是背景铺垫，与本次结论无直接关系）",
    "群里有人问能不能远程，这块按现有制度走，不在本议题。",
    "插一句无关的：饮水机换了供应商，口感大家自行适应。",
]


def slug(text: str) -> str:
    text = re.sub(r"[\\/:*?\"<>|\s]+", "_", text.strip())
    return text.strip("_")[:40]


def render_policy_md(d: Dossier, rnd: random.Random) -> str:
    """制度(md)：核心事实埋进条款散文，无 `核心事实：` 标签。"""
    noise = rnd.sample(NOISE_SENTENCES, 2)
    return f"""{DISCLAIMER}

# {COMPANY}{d.cat}内部制度：{d.proj}（{d.code}）

发文：{d.dept}　责任人：{d.owner}

第一条　适用范围。本制度仅适用于 {d.proj} 项目（代号 {d.code}）相关团队，外部通用做法不作为本项目执行依据。{noise[0]}

第二条　核心边界。{d.anchor}相关团队在答复外部咨询时，应以本条为准，不得用行业常识替代。

第三条　配套要求。{d.workflow}各环节负责人需留痕。{noise[1]}

第四条　历史参数。{d.temporal_subj}早期暂按 {d.temporal_old} 执行（注：本条后经复盘修订，最新值以复盘报告为准）。

第五条　解释权。本制度由{d.owner}（{d.dept}）负责解释，与本制度冲突的旧口径以最新复盘为准。

附则：本项目特别提示——{d.anti}
"""


def render_email_txt(d: Dossier, rnd: random.Random) -> str:
    """邮件线程(txt)：特殊例外藏在某封回复中间，带引用历史噪声。"""
    return f"""{DISCLAIMER}

发件人：{d.owner} <{d.code.lower()}@xingheng.example>
收件人：{d.dept}全体
主题：Re: Re: 关于{d.proj}（{d.code}）的几个执行口径

各位：

前面那封我快速回一下。第一个问题先不展开，按现有制度走。

关于例外情况，这里明确一下：{d.exception}这条是经过确认的，请按此执行，不要自行扩大范围。

> 在 上一封 中，有同事写道：
> 　这块到底有没有例外？能不能放宽？
> 　另外那个无关的会议室预订问题谁跟一下。

其它细节线下再约。

{d.owner}
{d.dept}
"""


def render_chat_txt(d: Dossier, rnd: random.Random) -> str:
    """群聊(txt)：结论数字从多条带时间戳消息里浮现，含口语/错别字。"""
    return f"""{DISCLAIMER}
[{d.proj} {d.code} 项目群 · 聊天记录导出]

09:02 {d.owner}：早，今天把 {d.chat_topic} 这事敲定下哈
09:03 小航：我这边有个数但不太确定，怕记串了
09:05 {d.owner}：别慌，咱以制度为准
09:11 阿茉：我翻了下纪要，应该是这个范围
09:14 {d.owner}：对，最终就按 {d.chat_value}，别再改了
09:15 小航：收到，那我同步给下游
09:16 {d.owner}：嗯，{d.chat_topic}就定 {d.chat_value}，谢谢各位
09:20 （有人发了个无关的表情包，略）
"""


def render_minutes_docx_blocks(d: Dossier, rnd: random.Random) -> list[tuple[str, str]]:
    """会议纪要(docx)：执行流程夹在跑题与待办之间。返回 (style, text) 块列表。"""
    noise = rnd.sample(NOISE_SENTENCES, 2)
    return [
        ("h0", f"{COMPANY}{d.proj}（{d.code}）项目纪要"),
        ("p", DISCLAIMER),
        ("p", f"主持：{d.owner}　记录：{d.dept}"),
        ("h2", "一、开场与跑题"),
        ("p", noise[0]),
        ("h2", "二、执行口径（重点）"),
        ("p", f"经讨论确认：{d.workflow}该流程为本项目强制要求，后续答复以此为准。"),
        ("p", noise[1]),
        ("h2", "三、待办"),
        ("bullet", f"{d.owner}：把上述流程同步到 {d.dept} wiki"),
        ("bullet", "其他：团建预算下次再议"),
    ]


def render_postmortem_pdf_blocks(d: Dossier) -> list[tuple[str, str]]:
    """复盘/事故报告(pdf)：把旧值修正为现值（时效冲突来源，正确答案=现值）。"""
    return [
        ("title", f"{COMPANY}{d.proj}（{d.code}）复盘报告"),
        ("h", "一、背景"),
        ("p", f"本次复盘聚焦 {d.temporal_subj}。该口径此前在早期文档中写为 {d.temporal_old}，实际运行后发现不合适。"),
        ("h", "二、修正结论"),
        ("p", f"经{d.owner}与{d.dept}确认，{d.temporal_subj}由 {d.temporal_old} 正式修正为 {d.temporal_new}，以现值为准。"),
        ("p", f"提醒：任何早期资料中出现的 {d.temporal_old} 均已作废，现行值为 {d.temporal_new}。"),
        ("h", "三、反常识提示"),
        ("p", d.anti),
    ]


def build_param_rows(d: Dossier) -> list[list[str]]:
    """参数表(xlsx)：把关键参数放到具体单元格，答案=某行某列。"""
    return [
        ["参数表", f"{d.proj}（{d.code}）", "责任人", d.owner],
        ["序号", "参数项", "数值", "备注"],
        ["1", "项目代号", d.code, d.dept],
        ["2", d.table_param, d.table_value, "本项目口径"],
        ["3", "评审顺序", "见会议纪要", "强制"],
        ["4", "数据级别", "内部资料", "仅本地检索评测"],
    ]


def build_slides(d: Dossier) -> list[tuple[str, list[str]]]:
    """幻灯片(pptx)：负责人 + 关键数字拆在标题与正文。"""
    return [
        (f"{d.proj}（{d.code}）项目概览",
         [DISCLAIMER, f"负责人：{d.owner}", f"归属：{d.dept}"]),
        (f"{d.proj} 关键口径",
         [d.slide_metric, f"代号 {d.code} 的口径以内部资料为准", d.anti]),
    ]


# ---------------------------------------------------------------------------
# 各格式落盘
# ---------------------------------------------------------------------------


def save_md(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def save_txt(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def save_docx_blocks(path: Path, blocks: list[tuple[str, str]]) -> None:
    doc = Document()
    doc.styles["Normal"].font.name = "Microsoft YaHei"
    doc.styles["Normal"].font.size = Pt(10.5)
    for style, text in blocks:
        if style == "h0":
            doc.add_heading(text, level=0)
        elif style == "h2":
            doc.add_heading(text, level=2)
        elif style == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        else:
            doc.add_paragraph(text)
    doc.save(path)


def save_pdf_blocks(path: Path, blocks: list[tuple[str, str]]) -> None:
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        font_name = "STSong-Light"
    except Exception:
        font_name = "Helvetica"
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("T", parent=styles["Title"], fontName=font_name,
                                 fontSize=15, leading=20, textColor=colors.HexColor("#1f2937"))
    body_style = ParagraphStyle("B", parent=styles["BodyText"], fontName=font_name,
                                fontSize=10, leading=15, spaceAfter=6)
    head_style = ParagraphStyle("H", parent=body_style, fontSize=11.5,
                                textColor=colors.HexColor("#0f766e"), spaceBefore=8, spaceAfter=4)
    doc = SimpleDocTemplate(str(path), pagesize=A4, rightMargin=42, leftMargin=42,
                            topMargin=42, bottomMargin=36)
    story = []
    for kind, text in blocks:
        if kind == "title":
            story.append(Paragraph(text, title_style))
            story.append(Spacer(1, 8))
        elif kind == "h":
            story.append(Paragraph(text, head_style))
        else:
            esc = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(esc, body_style))
    doc.build(story)


def save_xlsx(path: Path, rows: list[list[str]]) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "参数表"
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx, value=val)
            if r_idx <= 2:
                cell.font = Font(bold=True)
    wb.save(path)


def save_pptx(path: Path, slides: list[tuple[str, list[str]]]) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for line in bullets[1:]:
            p = body.add_paragraph()
            p.text = line
            p.font.size = PPt(16)
    prs.save(path)


# ---------------------------------------------------------------------------
# 干扰文档（220 份，不含任何信号答案）
# ---------------------------------------------------------------------------

NOISE_TOPICS = [
    ("行政通知", "关于办公区门禁与工位调整的说明"),
    ("培训材料", "新人入职安全意识培训讲义"),
    ("周报汇总", "本周各组进展与风险同步"),
    ("流程说明", "报销与差旅申请操作指引"),
    ("技术分享", "一次内部 RAG 检索调优经验分享"),
    ("公告", "节假日值班与系统冻结安排"),
    ("调研笔记", "竞品功能走查与可借鉴点"),
    ("术语表", "团队常用缩写与黑话对照"),
    ("环境说明", "测试与预发环境的使用约定"),
    ("会议室", "各楼层会议室预订规则"),
]
NOISE_BODY = [
    "本文为通用说明，不涉及任何具体项目的对外承诺或内部阈值。",
    "如与具体项目制度冲突，以对应项目的正式资料为准。",
    "以下内容仅作背景参考，请勿据此回答任何项目的事实性问题。",
    "相关同事如有疑问，请走对应流程提单，不在本文展开。",
    "本文档定期复审，过期内容以最新版本为准。",
]


def render_noise(idx: int, topic: tuple[str, str], rnd: random.Random,
                 decoy: str | None) -> tuple[str, list[tuple[str, str]], list[list[str]], list[tuple[str, list[str]]]]:
    """返回 (纯文本, docx块, xlsx行, pptx页)，调用方按格式取用其一。"""
    cat, title = topic
    lines = [DISCLAIMER, f"# {COMPANY}{cat}：{title}（NB-{idx:03d}）", ""]
    for _ in range(rnd.randint(3, 5)):
        lines.append(rnd.choice(NOISE_BODY))
    if decoy:
        # 近似代号干扰：提到一个与某信号项目相近、但事实完全不同的代号。
        lines.append(f"备注：本文提到的 {decoy} 与其它项目无关，请勿混淆，其数值均为占位示意。")
    text = "\n".join(lines)
    docx_blocks = [("h0", f"{COMPANY}{cat}：{title}")] + [("p", l) for l in lines[2:] if l]
    xlsx_rows = [["通用说明", title], ["项", "值"]] + [[str(i + 1), rnd.choice(NOISE_BODY)] for i in range(4)]
    pptx_slides = [(f"{cat}：{title}", [DISCLAIMER] + rnd.sample(NOISE_BODY, 2))]
    return text, docx_blocks, xlsx_rows, pptx_slides


# ---------------------------------------------------------------------------
# 特殊样本：长文 / 内嵌图片 / 扫描件（专门用来暴露"长文图谱开销"和"图片内容丢失"）
# ---------------------------------------------------------------------------

CJK_FONT_CANDIDATES = [
    r"C:\Windows\Fonts\msyh.ttc",
    r"C:\Windows\Fonts\msyh.ttf",
    r"C:\Windows\Fonts\simhei.ttf",
    r"C:\Windows\Fonts\simsun.ttc",
]


def _cjk_font(size: int):
    for fp in CJK_FONT_CANDIDATES:
        if os.path.exists(fp):
            try:
                return ImageFont.truetype(fp, size)
            except Exception:
                continue
    return ImageFont.load_default()


def render_text_image(lines: list[str], out_png: Path, size=(900, 360)) -> None:
    """把若干行文字渲染成 PNG 图片（文字只存在于像素里，文本抽取拿不到）。"""
    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)
    font = _cjk_font(30)
    y = 30
    for line in lines:
        draw.text((40, y), line, fill=(20, 30, 50), font=font)
        y += 52
    img.save(out_png)


def render_long_text(project: str, owner: str, buried_fact: str, target_chars: int,
                     rnd: random.Random) -> str:
    """生成一篇很长（target_chars 量级）的内部长文，把唯一事实埋在约 70% 处。"""
    filler = [
        f"{project}的执行细则在不同阶段反复迭代，本节复述背景与边界，供新成员通读。",
        "以下段落均为流程性铺陈，不含对外承诺，亦不构成任何阈值口径。",
        f"负责人{owner}多次强调，回答相关问题必须以内部资料为准，不得套用行业通用经验。",
        "本章节包含大量过程性描述、历史决议回顾、以及与本议题关联较弱的补充说明。",
        "为模拟真实长文档，这里刻意加入冗长的过渡性叙述，使关键事实被淹没在篇幅之中。",
    ]
    parts = [f"# {COMPANY}内部长文：{project} 操作汇编", DISCLAIMER, ""]
    body_len = 0
    inserted = False
    para_idx = 0
    while body_len < target_chars:
        para_idx += 1
        # 约 70% 处埋入唯一事实
        if not inserted and body_len >= target_chars * 0.7:
            parts.append(f"第 {para_idx} 条（关键）：{buried_fact}")
            inserted = True
        else:
            parts.append(f"第 {para_idx} 条：{rnd.choice(filler)}{rnd.choice(filler)}")
        body_len += len(parts[-1])
    if not inserted:
        parts.append(f"第 {para_idx + 1} 条（关键）：{buried_fact}")
    return "\n\n".join(parts)


def save_scanned_pdf(path: Path, image_png: Path) -> None:
    """把图片铺进 PDF，不写任何文本层 —— 模拟扫描件（lopdf 抽不到字）。"""
    c = pdf_canvas.Canvas(str(path), pagesize=A4)
    c.drawImage(str(image_png), 48, 360, width=500, preserveAspectRatio=True, mask="auto")
    c.save()


def save_image_docx(path: Path, paragraphs: list[str], image_png: Path) -> None:
    """docx：若干段正文（可抽取）+ 一张内嵌图片（图内文字抽取不到）。"""
    doc = Document()
    doc.styles["Normal"].font.name = "Microsoft YaHei"
    doc.styles["Normal"].font.size = Pt(10.5)
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.add_picture(str(image_png), width=DocxInches(5.5))
    doc.save(path)


def save_image_pptx(path: Path, title: str, bullets: list[str], image_png: Path) -> None:
    """pptx：标题+正文（可抽取）+ 一张图片（图内文字抽取不到）。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        tf.add_paragraph().text = b
    slide.shapes.add_picture(str(image_png), PptxInches(5.5), PptxInches(1.5), width=PptxInches(3.5))
    prs.save(path)


def build_specials(rnd: random.Random) -> list[dict]:
    """生成特殊样本并返回 manifest 记录（含 extractable 标记）。"""
    specials = []
    tmp = OUTPUT_DIR / "_img_tmp"
    tmp.mkdir(exist_ok=True)

    # 1) 长文 md（约 3 万字），事实埋在深处、可抽取
    p1 = OUTPUT_DIR / "special_001_长文_鸿图长卷_汇编.md"
    save_md(p1, render_long_text("鸿图长卷", "穆鸿", "鸿图长卷的归档密钥固定每 19 天轮换一次。", 30000, rnd))
    specials.append({"name": "鸿图长卷", "carrier": rel(p1), "format": "md", "kind": "long",
                     "facts": [{"hint": "归档密钥轮换", "clue": "每 19 天轮换一次", "extractable": True}]})

    # 2) 长文 docx（约 2.5 万字），可抽取
    p2 = OUTPUT_DIR / "special_002_长文_云梯手册_细则.docx"
    long2 = render_long_text("云梯手册", "云岫", "云梯手册规定客诉分级阈值固定为 23 分。", 25000, rnd)
    save_docx_blocks(p2, [("p", para) for para in long2.split("\n\n")])
    specials.append({"name": "云梯手册", "carrier": rel(p2), "format": "docx", "kind": "long",
                     "facts": [{"hint": "客诉分级阈值", "clue": "客诉分级阈值固定为 23 分", "extractable": True}]})

    # 3) 图文 docx：caption 事实可抽取，图内事实抽不到
    img3 = tmp / "img3.png"
    render_text_image(["晨曦项目 · 内部图表", "回滚阈值是 0.9%（图内文字，抽取不到）", "仅本地检索评测"], img3)
    p3 = OUTPUT_DIR / "special_003_图文_晨曦_带图.docx"
    save_image_docx(p3, [
        DISCLAIMER,
        "图1说明：晨曦项目的灰度比例固定为 6.3%。",  # caption fact, extractable
        "下图为晨曦项目的内部图表（关键阈值画在图里）。",
    ], img3)
    specials.append({"name": "晨曦项目", "carrier": rel(p3), "format": "docx", "kind": "image",
                     "facts": [
                         {"hint": "灰度比例", "clue": "灰度比例固定为 6.3%", "extractable": True},
                         {"hint": "回滚阈值", "clue": "回滚阈值是 0.9%", "extractable": False},
                     ]})

    # 4) 纯图 docx：关键事实只在图片里 → 抽不到
    img4 = tmp / "img4.png"
    render_text_image(["暮山项目 · 密钥规格", "暮山项目的密钥长度为 47 位", "（此页关键信息仅以图片呈现）"], img4)
    p4 = OUTPUT_DIR / "special_004_纯图_暮山_规格.docx"
    save_image_docx(p4, [DISCLAIMER, "本页关键规格以下图为准。"], img4)
    specials.append({"name": "暮山项目", "carrier": rel(p4), "format": "docx", "kind": "image_only",
                     "facts": [{"hint": "密钥长度", "clue": "密钥长度为 47 位", "extractable": False}]})

    # 5) 扫描件 PDF：整页是图，无文本层 → 抽不到
    img5 = tmp / "img5.png"
    render_text_image(["苍岭项目 内部扫描件", "苍岭项目的对账窗口为每月 8 号", "扫描自纸质文件，无文本层"], img5, size=(1000, 420))
    p5 = OUTPUT_DIR / "special_005_扫描件_苍岭_对账.pdf"
    save_scanned_pdf(p5, img5)
    specials.append({"name": "苍岭项目", "carrier": rel(p5), "format": "pdf", "kind": "scanned",
                     "facts": [{"hint": "对账窗口", "clue": "对账窗口为每月 8 号", "extractable": False}]})

    # 6) 图文 pptx：标题/正文事实可抽取，图内事实抽不到
    img6 = tmp / "img6.png"
    render_text_image(["白川项目 · 预算图", "白川项目的预算上限是 88 万（图内）"], img6)
    p6 = OUTPUT_DIR / "special_006_图文_白川_幻灯片.pptx"
    save_image_pptx(p6, "白川项目概览", [DISCLAIMER, "白川项目负责人是文澜。", "预算明细见右图。"], img6)
    specials.append({"name": "白川项目", "carrier": rel(p6), "format": "pptx", "kind": "image",
                     "facts": [
                         {"hint": "负责人", "clue": "白川项目负责人是文澜", "extractable": True},
                         {"hint": "预算上限", "clue": "预算上限是 88 万", "extractable": False},
                     ]})

    # 清理临时图片
    for f in tmp.glob("*.png"):
        f.unlink()
    tmp.rmdir()
    return specials


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------

# 每个信号项目固定 7 份体裁；格式映射如下（含 xlsx/pptx）。
SIGNAL_GENRES = [
    ("制度", "md"),
    ("邮件", "txt"),
    ("群聊", "txt"),
    ("纪要", "docx"),
    ("复盘", "pdf"),
    ("参数表", "xlsx"),
    ("幻灯片", "pptx"),
]


def clean_old() -> None:
    if OUTPUT_DIR.exists():
        for p in OUTPUT_DIR.glob("*"):
            if p.is_file() and (p.name.startswith(("sig_", "noise_", "special_")) or p.name == "corpus_manifest.json"):
                p.unlink()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def rel(path: Path) -> str:
    return f"Memory_Test_V2/{path.name}"


def main() -> int:
    rnd = random.Random(SEED)
    clean_old()

    manifest_dossiers = []
    serial = 0
    fmt_counter: dict[str, int] = {}

    for d in DOSSIERS:
        serial_base = serial
        carriers: dict[str, str] = {}
        for genre, ext in SIGNAL_GENRES:
            serial += 1
            fmt_counter[ext] = fmt_counter.get(ext, 0) + 1
            fname = f"sig_{serial:03d}_{slug(d.cat)}_{slug(d.proj)}_{genre}.{ext}"
            path = OUTPUT_DIR / fname
            if genre == "制度":
                save_md(path, render_policy_md(d, rnd))
            elif genre == "邮件":
                save_txt(path, render_email_txt(d, rnd))
            elif genre == "群聊":
                save_txt(path, render_chat_txt(d, rnd))
            elif genre == "纪要":
                save_docx_blocks(path, render_minutes_docx_blocks(d, rnd))
            elif genre == "复盘":
                save_pdf_blocks(path, render_postmortem_pdf_blocks(d))
            elif genre == "参数表":
                save_xlsx(path, build_param_rows(d))
            elif genre == "幻灯片":
                save_pptx(path, build_slides(d))
            carriers[genre] = rel(path)

        # 记录每条事实的承载文档/格式/精确 clue（供套件生成器消费）。
        manifest_dossiers.append({
            "code": d.code,
            "project": d.proj,
            "category": d.cat,
            "owner": d.owner,
            "department": d.dept,
            "decoy_code": d.decoy_code,
            "anti_common": d.anti,
            "facts": {
                "anchor": {"hint": "核心事实", "clue": d.anchor_clue,
                           "carrier": carriers["制度"], "format": "md", "genre": "制度"},
                "table": {"hint": d.table_param, "clue": d.table_value,
                          "carrier": carriers["参数表"], "format": "xlsx", "genre": "参数表"},
                "workflow": {"hint": "执行流程", "clue": d.workflow_clue,
                             "carrier": carriers["纪要"], "format": "docx", "genre": "纪要"},
                "exception": {"hint": "特殊例外", "clue": d.exception_clue,
                              "carrier": carriers["邮件"], "format": "txt", "genre": "邮件"},
                "chat": {"hint": d.chat_topic, "clue": d.chat_value,
                         "carrier": carriers["群聊"], "format": "txt", "genre": "群聊"},
                "temporal": {"hint": d.temporal_subj, "old": d.temporal_old, "clue": d.temporal_new,
                             "carrier": carriers["复盘"], "format": "pdf", "genre": "复盘"},
                "slide": {"hint": "关键数字", "clue": d.slide_clue,
                          "carrier": carriers["幻灯片"], "format": "pptx", "genre": "幻灯片"},
            },
            "signal_serials": [serial_base + 1, serial],
        })

    signal_count = serial

    # ---- 干扰文档 ----
    decoys = [d.decoy_code for d in DOSSIERS]
    noise_plan = [
        ("md", 40), ("txt", 40), ("docx", 35), ("pdf", 35), ("xlsx", 35), ("pptx", 35),
    ]
    nidx = 0
    for ext, count in noise_plan:
        for _ in range(count):
            nidx += 1
            fmt_counter[ext] = fmt_counter.get(ext, 0) + 1
            topic = rnd.choice(NOISE_TOPICS)
            decoy = decoys[nidx % len(decoys)] if nidx % 3 == 0 else None
            text, docx_blocks, xlsx_rows, pptx_slides = render_noise(nidx, topic, rnd, decoy)
            fname = f"noise_{nidx:03d}_{slug(topic[0])}.{ext}"
            path = OUTPUT_DIR / fname
            if ext == "md":
                save_md(path, text)
            elif ext == "txt":
                save_txt(path, text)
            elif ext == "docx":
                save_docx_blocks(path, docx_blocks)
            elif ext == "pdf":
                save_pdf_blocks(path, [("title", topic[1])] + [("p", l) for l in text.splitlines() if l][:6])
            elif ext == "xlsx":
                save_xlsx(path, xlsx_rows)
            elif ext == "pptx":
                save_pptx(path, pptx_slides)

    noise_count = nidx

    # ---- 特殊样本（长文 / 图片 / 扫描件）----
    specials = build_specials(rnd)
    special_count = len(specials)
    for s in specials:
        fmt_counter[s["format"]] = fmt_counter.get(s["format"], 0) + 1

    total = signal_count + noise_count + special_count

    manifest = {
        "company": COMPANY,
        "seed": SEED,
        "signal_count": signal_count,
        "noise_count": noise_count,
        "special_count": special_count,
        "total": total,
        "format_counts": fmt_counter,
        "signal_genres": SIGNAL_GENRES,
        "dossiers": manifest_dossiers,
        "specials": specials,
    }
    (OUTPUT_DIR / "corpus_manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"信号 {signal_count} + 干扰 {noise_count} + 特殊 {special_count} = {total} 份，输出到 {OUTPUT_DIR}")
    print("格式分布：", fmt_counter)
    print("manifest：", OUTPUT_DIR / "corpus_manifest.json")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
